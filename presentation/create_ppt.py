import os
import sys

def main():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        print("python-pptx is not installed. Installing python-pptx...")
        import subprocess
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Constants
    BG_COLOR = RGBColor(5, 8, 20)        # Deep Cyberpunk Blue (#050814)
    CYAN = RGBColor(0, 240, 255)         # Neon Cyan (#00f0ff)
    MAGENTA = RGBColor(255, 0, 127)      # Neon Magenta (#ff007f)
    GOLD = RGBColor(255, 170, 0)         # Neon Gold (#ffaa00)
    WHITE = RGBColor(240, 240, 245)      # Light Grey/White (#f0f0f5)
    MUTED = RGBColor(148, 163, 184)      # Slate Muted (#94a3b8)
    RED = RGBColor(255, 51, 68)          # Warning Red (#ff3344)

    def apply_dark_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_corner_borders(slide):
        # Draw top and bottom thin accent lines to look like HUD
        shapes = slide.shapes
        # Top Accent Line
        top_line = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.8), Inches(12.333), Inches(0.02))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = CYAN
        top_line.line.color.rgb = CYAN

        # Bottom Accent Line
        bottom_line = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.02))
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = CYAN
        bottom_line.line.color.rgb = CYAN

        # HUD Text on Top
        tx_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(5.0), Inches(0.4))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = "OMS // AI SURVEILLANCE & REAL-TIME MONITORING"
        p.font.name = "Courier New"
        p.font.size = Pt(11)
        p.font.color.rgb = CYAN
        p.font.bold = True

        tx_box_status = shapes.add_textbox(Inches(9.333), Inches(0.3), Inches(3.5), Inches(0.4))
        tf_status = tx_box_status.text_frame
        p_status = tf_status.paragraphs[0]
        p_status.alignment = PP_ALIGN.RIGHT
        p_status.text = "SYS_STATUS: ACTIVE"
        p_status.font.name = "Courier New"
        p_status.font.size = Pt(11)
        p_status.font.color.rgb = RGBColor(0, 255, 102)
        p_status.font.bold = True

    def add_slide_header(slide, title_text, slide_num):
        shapes = slide.shapes
        tx_box = shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.333), Inches(0.8))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        
        # Add Slide Number
        run_num = p.add_run()
        run_num.text = f"{slide_num:02d} "
        run_num.font.name = "Courier New"
        run_num.font.size = Pt(24)
        run_num.font.color.rgb = MAGENTA
        run_num.font.bold = True
        
        # Add Title
        run_title = p.add_run()
        run_title.text = title_text.upper()
        run_title.font.name = "Arial"
        run_title.font.size = Pt(28)
        run_title.font.color.rgb = WHITE
        run_title.font.bold = True

        # Accent thin underline for title
        underline = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(4.5), Inches(0.04))
        underline.fill.solid()
        underline.fill.fore_color.rgb = MAGENTA
        underline.line.color.rgb = MAGENTA

    # --- SLIDE 1: Title Slide ---
    slide_layout = prs.slide_layouts[6] # Blank Layout
    slide1 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide1)

    # Frame Box for Title
    frame = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(3.2))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(10, 18, 42)
    frame.line.color.rgb = CYAN
    frame.line.width = Pt(2)

    # Main Title
    tx_title = slide1.shapes.add_textbox(Inches(1.5), Inches(1.4), Inches(10.333), Inches(1.2))
    p = tx_title.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "OMS"
    run.font.name = "Arial"
    run.font.size = Pt(64)
    run.font.color.rgb = WHITE
    run.font.bold = True

    # Subtitle
    tx_sub = slide1.shapes.add_textbox(Inches(1.5), Inches(2.6), Inches(10.333), Inches(1.0))
    tf_sub = tx_sub.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.alignment = PP_ALIGN.CENTER
    run_sub = p_sub.add_run()
    run_sub.text = "OBJECT MONITORING SYSTEM"
    run_sub.font.name = "Arial"
    run_sub.font.size = Pt(22)
    run_sub.font.color.rgb = MAGENTA
    run_sub.font.bold = True
    
    p_tag = tf_sub.add_paragraph()
    p_tag.alignment = PP_ALIGN.CENTER
    run_tag = p_tag.add_run()
    run_tag.text = "AI-Powered Smart Surveillance & Real-Time Monitoring Platform"
    run_tag.font.name = "Courier New"
    run_tag.font.size = Pt(13)
    run_tag.font.color.rgb = MUTED

    # Overview Box
    overview_box = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.7), Inches(5.0), Inches(2.0))
    overview_box.fill.solid()
    overview_box.fill.fore_color.rgb = RGBColor(8, 25, 48)
    overview_box.line.color.rgb = CYAN
    tf_ov = overview_box.text_frame
    tf_ov.word_wrap = True
    p_ov = tf_ov.paragraphs[0]
    run_ov_lbl = p_ov.add_run()
    run_ov_lbl.text = "[OVERVIEW]\n"
    run_ov_lbl.font.name = "Courier New"
    run_ov_lbl.font.size = Pt(12)
    run_ov_lbl.font.color.rgb = CYAN
    run_ov_lbl.font.bold = True
    
    run_ov_txt = p_ov.add_run()
    run_ov_txt.text = "OMS is an AI-based surveillance system that monitors live camera feeds, detects people/objects in real time, tracks activity, and features an interactive PyQt6 desktop dashboard."
    run_ov_txt.font.name = "Arial"
    run_ov_txt.font.size = Pt(12)
    run_ov_txt.font.color.rgb = WHITE

    # Project Details
    details_box = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.833), Inches(4.7), Inches(5.0), Inches(2.0))
    details_box.fill.solid()
    details_box.fill.fore_color.rgb = RGBColor(15, 10, 25)
    details_box.line.color.rgb = MAGENTA
    tf_det = details_box.text_frame
    tf_det.word_wrap = True
    p_det = tf_det.paragraphs[0]
    run_det_lbl = p_det.add_run()
    run_det_lbl.text = "[ACADEMIC DETAILS]\n"
    run_det_lbl.font.name = "Courier New"
    run_det_lbl.font.size = Pt(12)
    run_det_lbl.font.color.rgb = MAGENTA
    run_det_lbl.font.bold = True
    
    run_det_txt = p_det.add_run()
    run_det_txt.text = "Developer: Student Name\nDept: Computer Science & Engineering\nYear: 2025 - 2026\nRepo: github.com/yourprofile/oms"
    run_det_txt.font.name = "Arial"
    run_det_txt.font.size = Pt(12)
    run_det_txt.font.color.rgb = WHITE

    # --- SLIDE 2: Problem Statement ---
    slide2 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide2)
    add_corner_borders(slide2)
    add_slide_header(slide2, "Problem Statement", 1)

    problems = [
        ("Manual Fatigue", "Traditional CCTV systems require continuous, passive human monitoring.", "⚠️"),
        ("Human Error", "Operators frequently miss crucial, sudden, or stealthy security threats.", "📉"),
        ("Multi-Feed Overhead", "Actively monitoring multiple camera feeds simultaneously is highly inefficient.", "🖥️"),
        ("No Scene Analysis", "Standard surveillance networks cannot classify items or detect intrusions automatically.", "🔒")
    ]

    card_width = Inches(2.8)
    card_height = Inches(2.5)
    spacing = Inches(0.3)
    start_left = Inches(0.5)
    start_top = Inches(2.2)

    for i, (title, desc, icon) in enumerate(problems):
        left = start_left + i * (card_width + spacing)
        card = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, start_top, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(10, 15, 30)
        card.line.color.rgb = RED
        card.line.width = Pt(1.5)
        
        tf = card.text_frame
        tf.word_wrap = True
        p_icon = tf.paragraphs[0]
        run_i = p_icon.add_run()
        run_i.text = f"{icon}\n\n"
        run_i.font.size = Pt(20)
        
        p_t = tf.add_paragraph()
        run_t = p_t.add_run()
        run_t.text = f"{title}\n"
        run_t.font.name = "Arial"
        run_t.font.size = Pt(14)
        run_t.font.color.rgb = WHITE
        run_t.font.bold = True
        
        p_d = tf.add_paragraph()
        run_d = p_d.add_run()
        run_d.text = desc
        run_d.font.name = "Arial"
        run_d.font.size = Pt(11)
        run_d.font.color.rgb = MUTED

    # CCTV vs OMS block
    vs_box = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.0), Inches(12.333), Inches(1.3))
    vs_box.fill.solid()
    vs_box.fill.fore_color.rgb = RGBColor(5, 5, 10)
    vs_box.line.color.rgb = CYAN
    tf_vs = vs_box.text_frame
    tf_vs.word_wrap = True
    p_vs = tf_vs.paragraphs[0]
    p_vs.alignment = PP_ALIGN.CENTER
    run_vs1 = p_vs.add_run()
    run_vs1.text = "TRADITIONAL CCTV: Passive Recording & Manual Review Only\n"
    run_vs1.font.name = "Arial"
    run_vs1.font.size = Pt(13)
    run_vs1.font.color.rgb = RED
    run_vs1.font.bold = True
    
    run_vs_mid = p_vs.add_run()
    run_vs_mid.text = "  VS  \n"
    run_vs_mid.font.name = "Courier New"
    run_vs_mid.font.size = Pt(14)
    run_vs_mid.font.color.rgb = MAGENTA
    run_vs_mid.font.bold = True

    run_vs2 = p_vs.add_run()
    run_vs2.text = "OMS SMART SURVEILLANCE: Automated Real-Time AI Scene Analysis"
    run_vs2.font.name = "Arial"
    run_vs2.font.size = Pt(13)
    run_vs2.font.color.rgb = CYAN
    run_vs2.font.bold = True

    # --- SLIDE 3: Objectives ---
    slide3 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide3)
    add_corner_borders(slide3)
    add_slide_header(slide3, "Objectives of the Project", 2)

    # Objectives text boxes on the left
    tx_obj = slide3.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(7.5), Inches(4.3))
    tf_obj = tx_obj.text_frame
    tf_obj.word_wrap = True

    objectives = [
        ("Real-time Surveillance Feed Monitoring", "Process and analyze live frame rates automatically."),
        ("Object Detection using YOLO Engine", "Localize and classify subjects (people, luggage, cars) instantly."),
        ("Face Recognition Integrations", "Identify known employees, students, or visitors inside monitored zones."),
        ("Movement Tracking & History", "Track unique target trajectories and assign IDs across frames."),
        ("Automated Security Event Logging", "Log logs, timestamps, alerts, and classifications to disk database."),
        ("Fast Desktop PyQt6 Dashboard Console", "Present live visual tracking data and active logs grid to operators.")
    ]

    for title, desc in objectives:
        p = tf_obj.add_paragraph() if tf_obj.paragraphs[0].text else tf_obj.paragraphs[0]
        p.space_after = Pt(8)
        run_dot = p.add_run()
        run_dot.text = "» "
        run_dot.font.name = "Courier New"
        run_dot.font.size = Pt(14)
        run_dot.font.color.rgb = CYAN
        run_dot.font.bold = True

        run_t = p.add_run()
        run_t.text = f"{title}: "
        run_t.font.name = "Arial"
        run_t.font.size = Pt(13)
        run_t.font.color.rgb = WHITE
        run_t.font.bold = True

        run_d = p.add_run()
        run_d.text = desc
        run_d.font.name = "Arial"
        run_d.font.size = Pt(12)
        run_d.font.color.rgb = MUTED

    # Cyberpunk target radar box on right
    radar_box = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.0), Inches(2.5), Inches(3.2), Inches(3.2))
    radar_box.fill.solid()
    radar_box.fill.fore_color.rgb = RGBColor(3, 15, 30)
    radar_box.line.color.rgb = CYAN
    radar_box.line.width = Pt(2)
    tf_rad = radar_box.text_frame
    tf_rad.word_wrap = True
    p_rad = tf_rad.paragraphs[0]
    p_rad.alignment = PP_ALIGN.CENTER
    p_rad.space_before = Pt(80)
    run_rad = p_rad.add_run()
    run_rad.text = "OMS RADAR\n[TARGET_SCAN: OK]"
    run_rad.font.name = "Courier New"
    run_rad.font.size = Pt(12)
    run_rad.font.color.rgb = CYAN
    run_rad.font.bold = True

    # --- SLIDE 4: System Architecture ---
    slide4 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide4)
    add_corner_borders(slide4)
    add_slide_header(slide4, "System Architecture", 3)

    blocks = [
        ("Camera Feed", "Webcam/RTSP Input"),
        ("Frame Capture", "OpenCV Decoder"),
        ("YOLO Core", "AI Detector"),
        ("Tracking Engine", "Embedding Scan"),
        ("Event Logger", "Persistence File"),
        ("Live UI", "PyQt6 Dashboard")
    ]

    card_w = Inches(1.8)
    card_h = Inches(1.8)
    spacing_b = Inches(0.2)
    start_x = Inches(0.5)
    start_y = Inches(2.6)

    for i, (name, role) in enumerate(blocks):
        left_b = start_x + i * (card_w + spacing_b)
        block = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_b, start_y, card_w, card_h)
        block.fill.solid()
        
        # Color highlight for YOLO and Tracking modules
        if "YOLO" in name:
            block.fill.fore_color.rgb = RGBColor(12, 45, 64)
            block.line.color.rgb = CYAN
            block.line.width = Pt(2.5)
        elif "Tracking" in name:
            block.fill.fore_color.rgb = RGBColor(45, 12, 35)
            block.line.color.rgb = MAGENTA
            block.line.width = Pt(2.5)
        else:
            block.fill.fore_color.rgb = RGBColor(10, 15, 30)
            block.line.color.rgb = MUTED
            block.line.width = Pt(1)

        tf = block.text_frame
        tf.word_wrap = True
        p_name = tf.paragraphs[0]
        p_name.alignment = PP_ALIGN.CENTER
        p_name.space_after = Pt(5)
        run_name = p_name.add_run()
        run_name.text = f"{name}\n"
        run_name.font.name = "Arial"
        run_name.font.size = Pt(12)
        run_name.font.color.rgb = WHITE
        run_name.font.bold = True

        p_role = tf.add_paragraph()
        p_role.alignment = PP_ALIGN.CENTER
        run_role = p_role.add_run()
        run_role.text = role
        run_role.font.name = "Courier New"
        run_role.font.size = Pt(10)
        run_role.font.color.rgb = MUTED

        # Add connection indicator arrow except for the last block
        if i < len(blocks) - 1:
            tx_arrow = slide4.shapes.add_textbox(left_b + card_w, start_y + Inches(0.6), spacing_b, Inches(0.6))
            p_arr = tx_arrow.text_frame.paragraphs[0]
            p_arr.alignment = PP_ALIGN.CENTER
            run_arr = p_arr.add_run()
            run_arr.text = "➔"
            run_arr.font.size = Pt(14)
            run_arr.font.color.rgb = MAGENTA
            run_arr.font.bold = True

    # Description text
    desc_box = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.0), Inches(12.333), Inches(1.3))
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = RGBColor(5, 12, 28)
    desc_box.line.color.rgb = CYAN
    tf_ds = desc_box.text_frame
    tf_ds.word_wrap = True
    p_ds = tf_ds.paragraphs[0]
    run_ds_lbl = p_ds.add_run()
    run_ds_lbl.text = "[PIPELINE DESCRIPTION]\n"
    run_ds_lbl.font.name = "Courier New"
    run_ds_lbl.font.size = Pt(12)
    run_ds_lbl.font.color.rgb = CYAN
    run_ds_lbl.font.bold = True
    
    run_ds_txt = p_ds.add_run()
    run_ds_txt.text = "OMS receives video from a camera source, processes each frame, detects people/objects, performs recognition and tracking where needed, logs important events, and displays the results in the dashboard."
    run_ds_txt.font.name = "Arial"
    run_ds_txt.font.size = Pt(12)
    run_ds_txt.font.color.rgb = WHITE

    # --- SLIDE 5: Working & Modules ---
    slide5 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide5)
    add_corner_borders(slide5)
    add_slide_header(slide5, "Working of OMS / Main Modules", 4)

    # Columns: Left (Modules), Right (Workflow Loop)
    tx_mod = slide5.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.8), Inches(4.3))
    tf_mod = tx_mod.text_frame
    tf_mod.word_wrap = True
    p_m_head = tf_mod.paragraphs[0]
    p_m_head.space_after = Pt(10)
    run_mh = p_m_head.add_run()
    run_mh.text = "Core Modules:\n"
    run_mh.font.name = "Arial"
    run_mh.font.size = Pt(16)
    run_mh.font.color.rgb = CYAN
    run_mh.font.bold = True

    modules_info = [
        ("1. Camera Input Module", "Interfaces with usb cams, RTSP network feeds, and archives."),
        ("2. Object Detection Module", "Passes frames to YOLO to return boundary coordinates."),
        ("3. Face Recognition Module", "Compares facial crops with database models."),
        ("4. Tracking & Logging Engine", "Maintains target state and updates log registers.")
    ]

    for title, desc in modules_info:
        p = tf_mod.add_paragraph()
        p.space_after = Pt(8)
        run_mt = p.add_run()
        run_mt.text = f"{title}\n"
        run_mt.font.name = "Arial"
        run_mt.font.size = Pt(12)
        run_mt.font.color.rgb = WHITE
        run_mt.font.bold = True
        
        run_md = p.add_run()
        run_md.text = desc
        run_md.font.name = "Arial"
        run_md.font.size = Pt(11)
        run_md.font.color.rgb = MUTED

    # Right side: Workflow Loop
    tx_flow = slide5.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(6.0), Inches(4.3))
    tf_flow = tx_flow.text_frame
    tf_flow.word_wrap = True
    p_f_head = tf_flow.paragraphs[0]
    p_f_head.space_after = Pt(10)
    run_fh = p_f_head.add_run()
    run_fh.text = "Active Workflow Loop:\n"
    run_fh.font.name = "Arial"
    run_fh.font.size = Pt(16)
    run_fh.font.color.rgb = MAGENTA
    run_fh.font.bold = True

    workflow_steps = [
        ("01", "Stream Capture", "Continuously read and capture raw frame buffers."),
        ("02", "AI Model Analysis", "Feed frames to YOLO model for target identification."),
        ("03", "Trajectory Track", "Associate boundaries with IDs and scan registered faces."),
        ("04", "Database Persistence", "Commit alert updates to local disk log files."),
        ("05", "Refresh GUI Dashboard", "Render annotated frame streams and update logs lists.")
    ]

    for num, step_n, step_d in workflow_steps:
        p = tf_flow.add_paragraph()
        p.space_after = Pt(6)
        
        run_num = p.add_run()
        run_num.text = f"[{num}] "
        run_num.font.name = "Courier New"
        run_num.font.size = Pt(12)
        run_num.font.color.rgb = MAGENTA
        run_num.font.bold = True
        
        run_sn = p.add_run()
        run_sn.text = f"{step_n}: "
        run_sn.font.name = "Arial"
        run_sn.font.size = Pt(12)
        run_sn.font.color.rgb = WHITE
        run_sn.font.bold = True

        run_sd = p.add_run()
        run_sd.text = step_d
        run_sd.font.name = "Arial"
        run_sd.font.size = Pt(11)
        run_sd.font.color.rgb = MUTED

    # --- SLIDE 6: Features & Technologies ---
    slide6 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide6)
    add_corner_borders(slide6)
    add_slide_header(slide6, "Key Features & Technologies Used", 5)

    # Left: Features
    tx_feat = slide6.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.8), Inches(4.3))
    tf_feat = tx_feat.text_frame
    tf_feat.word_wrap = True
    p_ft_head = tf_feat.paragraphs[0]
    p_ft_head.space_after = Pt(12)
    run_fth = p_ft_head.add_run()
    run_fth.text = "Key Features:\n"
    run_fth.font.name = "Arial"
    run_fth.font.size = Pt(16)
    run_fth.font.color.rgb = CYAN
    run_fth.font.bold = True

    features_list = [
        "Multi-source streams compatibility (Webcams, CCTV IP RTSP feeds).",
        "High-performance detection with optimized inference latency.",
        "Entity tracking to avoid redundant logging of active targets.",
        "Persistent logs database with timestamps and confidence records.",
        "Interactive PyQt6 desktop control UI dashboard panel."
    ]

    for feat in features_list:
        p = tf_feat.add_paragraph()
        p.space_after = Pt(8)
        run_dot = p.add_run()
        run_dot.text = "» "
        run_dot.font.name = "Courier New"
        run_dot.font.size = Pt(13)
        run_dot.font.color.rgb = CYAN
        run_dot.font.bold = True
        
        run_txt = p.add_run()
        run_txt.text = feat
        run_txt.font.name = "Arial"
        run_txt.font.size = Pt(12)
        run_txt.font.color.rgb = WHITE

    # Right: Technologies Table
    tx_tab = slide6.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(6.0), Inches(4.3))
    tf_tab = tx_tab.text_frame
    tf_tab.word_wrap = True
    p_tb_head = tf_tab.paragraphs[0]
    p_tb_head.space_after = Pt(12)
    run_tbh = p_tb_head.add_run()
    run_tbh.text = "Technology Stack Purpose:\n"
    run_tbh.font.name = "Arial"
    run_tbh.font.size = Pt(16)
    run_tbh.font.color.rgb = MAGENTA
    run_tbh.font.bold = True

    tech_data = [
        ("Python", "Core system orchestration & scripting logic."),
        ("OpenCV", "Stream frame decoding & rendering processes."),
        ("YOLO Engine", "Deep learning model for multi-object detection."),
        ("PyQt6", "Responsive graphical user dashboard interface."),
        ("Face Recognition", "Generates facial embeddings and matching."),
        ("NumPy", "Efficient frame array calculations & manipulation.")
    ]

    for tech, purpose in tech_data:
        p = tf_tab.add_paragraph()
        p.space_after = Pt(6)
        run_tc = p.add_run()
        run_tc.text = f"{tech}: "
        run_tc.font.name = "Arial"
        run_tc.font.size = Pt(12)
        run_tc.font.color.rgb = CYAN
        run_tc.font.bold = True
        
        run_pur = p.add_run()
        run_pur.text = purpose
        run_pur.font.name = "Arial"
        run_pur.font.size = Pt(11)
        run_pur.font.color.rgb = MUTED

    # --- SLIDE 7: Applications & Scope ---
    slide7 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide7)
    add_corner_borders(slide7)
    add_slide_header(slide7, "Applications and Future Scope", 6)

    # Applications on left
    tx_app = slide7.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.8), Inches(4.3))
    tf_app = tx_app.text_frame
    tf_app.word_wrap = True
    p_ap_head = tf_app.paragraphs[0]
    p_ap_head.space_after = Pt(10)
    run_aph = p_ap_head.add_run()
    run_aph.text = "Target Deployment Environments:\n"
    run_aph.font.name = "Arial"
    run_aph.font.size = Pt(16)
    run_aph.font.color.rgb = CYAN
    run_aph.font.bold = True

    apps = [
        ("Computer Science Labs", "Log access controls of restricted computing areas."),
        ("Educational Campuses", "Monitor campus perimeters and identify visitors."),
        ("Corporate Offices", "Log worker attendance markers automatically."),
        ("Residential Security", "Enhance home access gates with automated alarms.")
    ]

    for title, desc in apps:
        p = tf_app.add_paragraph()
        p.space_after = Pt(8)
        run_lbl = p.add_run()
        run_lbl.text = f"• {title}: "
        run_lbl.font.name = "Arial"
        run_lbl.font.size = Pt(12)
        run_lbl.font.color.rgb = WHITE
        run_lbl.font.bold = True
        
        run_dsc = p.add_run()
        run_dsc.text = desc
        run_dsc.font.name = "Arial"
        run_dsc.font.size = Pt(11)
        run_dsc.font.color.rgb = MUTED

    # Future Scope on right
    tx_fut = slide7.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(6.0), Inches(4.3))
    tf_fut = tx_fut.text_frame
    tf_fut.word_wrap = True
    p_ft_head = tf_fut.paragraphs[0]
    p_ft_head.space_after = Pt(10)
    run_fth = p_ft_head.add_run()
    run_fth.text = "Future Expansion Roadmap:\n"
    run_fth.font.name = "Arial"
    run_fth.font.size = Pt(16)
    run_fth.font.color.rgb = MAGENTA
    run_fth.font.bold = True

    roadmap = [
        ("Behavior Anomaly Flags", "Trigger logs for loitering, falls, or fights."),
        ("Instant Mobile Alerts", "Send immediate threat snapshots via Telegram."),
        ("Multi-camera Central Nodes", "Scale pipeline load across central servers."),
        ("Cloud Analytics Sync", "Publish local logs to remote browser dashboards.")
    ]

    for title, desc in roadmap:
        p = tf_fut.add_paragraph()
        p.space_after = Pt(8)
        run_lbl = p.add_run()
        run_lbl.text = f"» {title}: "
        run_lbl.font.name = "Courier New"
        run_lbl.font.size = Pt(12)
        run_lbl.font.color.rgb = MAGENTA
        run_lbl.font.bold = True
        
        run_dsc = p.add_run()
        run_dsc.text = desc
        run_dsc.font.name = "Arial"
        run_dsc.font.size = Pt(11)
        run_dsc.font.color.rgb = MUTED

    # --- SLIDE 8: Conclusion ---
    slide8 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide8)
    add_corner_borders(slide8)
    add_slide_header(slide8, "Conclusion", 7)

    # Conclusion box in the center
    concl_box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.2), Inches(10.333), Inches(2.0))
    concl_box.fill.solid()
    concl_box.fill.fore_color.rgb = RGBColor(10, 18, 42)
    concl_box.line.color.rgb = MAGENTA
    concl_box.line.width = Pt(2)
    tf_co = concl_box.text_frame
    tf_co.word_wrap = True
    p_co = tf_co.paragraphs[0]
    p_co.alignment = PP_ALIGN.CENTER
    run_co_h = p_co.add_run()
    run_co_h.text = "SUMMARY & CONCLUSION\n\n"
    run_co_h.font.name = "Arial"
    run_co_h.font.size = Pt(15)
    run_co_h.font.color.rgb = MAGENTA
    run_co_h.font.bold = True

    run_co_t = p_co.add_run()
    run_co_t.text = "OMS is an AI-powered smart surveillance system that improves traditional CCTV monitoring by adding real-time object detection, tracking, recognition, and event logging. It helps reduce manual monitoring effort and provides a more intelligent and efficient security monitoring solution."
    run_co_t.font.name = "Arial"
    run_co_t.font.size = Pt(12)
    run_co_t.font.color.rgb = WHITE

    # GitHub Details
    gh_box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.5), Inches(5.0), Inches(1.5))
    gh_box.fill.solid()
    gh_box.fill.fore_color.rgb = RGBColor(5, 12, 28)
    gh_box.line.color.rgb = CYAN
    tf_gh = gh_box.text_frame
    tf_gh.word_wrap = True
    p_gh = tf_gh.paragraphs[0]
    run_gh_lbl = p_gh.add_run()
    run_gh_lbl.text = "[OMS GITHUB REPOSITORY]\n"
    run_gh_lbl.font.name = "Courier New"
    run_gh_lbl.font.size = Pt(11)
    run_gh_lbl.font.color.rgb = CYAN
    run_gh_lbl.font.bold = True
    
    run_gh_txt = p_gh.add_run()
    run_gh_txt.text = "Developer Profile: github.com/yourprofile\nRepository: github.com/yourprofile/oms"
    run_gh_txt.font.name = "Arial"
    run_gh_txt.font.size = Pt(11)
    run_gh_txt.font.color.rgb = WHITE

    # Thank You
    ty_box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.833), Inches(4.5), Inches(5.0), Inches(1.5))
    ty_box.fill.solid()
    ty_box.fill.fore_color.rgb = RGBColor(5, 5, 10)
    ty_box.line.color.rgb = CYAN
    tf_ty = ty_box.text_frame
    tf_ty.word_wrap = True
    p_ty = tf_ty.paragraphs[0]
    p_ty.alignment = PP_ALIGN.CENTER
    p_ty.space_before = Pt(15)
    run_ty = p_ty.add_run()
    run_ty.text = "THANK YOU\n"
    run_ty.font.name = "Arial"
    run_ty.font.size = Pt(28)
    run_ty.font.color.rgb = CYAN
    run_ty.font.bold = True
    
    run_qa = p_ty.add_run()
    run_qa.text = "Q&A Session // Questions Welcome"
    run_qa.font.name = "Courier New"
    run_qa.font.size = Pt(10)
    run_qa.font.color.rgb = MUTED

    prs.save("OMS_Presentation.pptx")
    print("PowerPoint presentation generated successfully as 'OMS_Presentation.pptx'")

if __name__ == "__main__":
    main()
